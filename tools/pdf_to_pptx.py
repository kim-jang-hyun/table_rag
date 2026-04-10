"""PDF → PPTX converter (real table objects).

Extracts tables from a PDF using PyMuPDF and renders them as actual PowerPoint
``Table`` shapes via python-pptx.  Pages without tables become text slides.

Features
--------
- Cross-page table merging (via ``table_rag.table``)
- rowspan back-fill and multi-row header merging (via ``table_rag.table``)
- One slide per table (title + table shape)
- Text-only pages rendered as text slides

Usage::

    python tools/pdf_to_pptx.py

Or call programmatically::

    from tools.pdf_to_pptx import pdf_to_pptx
    from pathlib import Path

    info = pdf_to_pptx(Path("report.pdf"), Path("report.pptx"))
    print(info)
"""

import sys
from pathlib import Path
from typing import List

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Pt

from table_rag.table import (
    extract_raw_tables_from_doc,
    find_table_title,
    merge_cross_page_raw_tables,
    normalize_table,
)

# ── Slide layout constants ─────────────────────────────────────────────────────

SLIDE_W = Emu(9_144_000)   # 10 inches
SLIDE_H = Emu(6_858_000)   # 7.5 inches

MARGIN  = Emu(457_200)     # 0.5 inch
TITLE_H = Emu(548_640)     # title box height
GAP     = Emu(182_880)     # gap between title and table
TABLE_Y = MARGIN + TITLE_H + GAP
TABLE_H = SLIDE_H - TABLE_Y - MARGIN
TABLE_W = SLIDE_W - MARGIN * 2

COLOR_HEADER_BG = RGBColor(0x2F, 0x54, 0x96)  # dark blue
COLOR_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)  # white
COLOR_ROW_ODD   = RGBColor(0xDF, 0xE8, 0xF5)  # light blue (odd rows)
COLOR_ROW_EVEN  = RGBColor(0xFF, 0xFF, 0xFF)  # white (even rows)
COLOR_CELL_TEXT = RGBColor(0x26, 0x26, 0x26)  # near-black

HEADER_FONT_PT = 11
BODY_FONT_PT   = 10


# ── PDF extraction ─────────────────────────────────────────────────────────────

def _extract_tables_from_pdf(pdf_path: Path) -> List[dict]:
    """Return merged table dicts from *pdf_path* (cross-page merging enabled)."""
    doc = fitz.open(str(pdf_path))
    raw = extract_raw_tables_from_doc(doc)
    doc.close()
    if not raw:
        return []
    return merge_cross_page_raw_tables(raw)


def _extract_text_by_page(pdf_path: Path) -> dict:
    """Return a mapping of {page_no: text} for all pages in *pdf_path*."""
    doc = fitz.open(str(pdf_path))
    result = {}
    for page_index in range(doc.page_count):
        page = doc.load_page(page_index)
        text = " ".join((page.get_text("text") or "").split())
        result[page_index + 1] = text
    doc.close()
    return result


# ── PPTX builders ─────────────────────────────────────────────────────────────

def _add_title_box(slide, title: str) -> None:
    txBox = slide.shapes.add_textbox(MARGIN, MARGIN, TABLE_W, TITLE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_CELL_TEXT


def _set_cell_style(cell, text: str, *, is_header: bool, row_idx: int) -> None:
    cell.text = text
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.font.size = Pt(HEADER_FONT_PT if is_header else BODY_FONT_PT)
    para.font.bold = is_header
    para.font.color.rgb = COLOR_HEADER_FG if is_header else COLOR_CELL_TEXT

    fill = cell.fill
    fill.solid()
    if is_header:
        fill.fore_color.rgb = COLOR_HEADER_BG
    elif row_idx % 2 == 0:
        fill.fore_color.rgb = COLOR_ROW_ODD
    else:
        fill.fore_color.rgb = COLOR_ROW_EVEN


def _add_table_slide(
    prs: Presentation, title: str, header: List[str], body: List[List[str]]
) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    if title:
        _add_title_box(slide, title)

    n_rows = 1 + len(body)
    n_cols = len(header)
    if n_rows < 1 or n_cols < 1:
        return

    col_w = TABLE_W // n_cols
    tbl_top = TABLE_Y if title else MARGIN
    tbl_h = SLIDE_H - tbl_top - MARGIN

    shape = slide.shapes.add_table(n_rows, n_cols, MARGIN, tbl_top, TABLE_W, tbl_h)
    tbl = shape.table

    for c in range(n_cols):
        tbl.columns[c].width = col_w

    for c, col_name in enumerate(header):
        _set_cell_style(tbl.cell(0, c), col_name, is_header=True, row_idx=0)

    for r, row in enumerate(body, start=1):
        for c, val in enumerate(row):
            _set_cell_style(tbl.cell(r, c), val, is_header=False, row_idx=r)


def _add_text_slide(prs: Presentation, page_no: int, text: str) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    _add_title_box(slide, f"페이지 {page_no} — 텍스트")

    txBox = slide.shapes.add_textbox(MARGIN, TABLE_Y, TABLE_W, TABLE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text[:2000]
    p.font.size = Pt(10)
    p.font.color.rgb = COLOR_CELL_TEXT


# ── Main converter ─────────────────────────────────────────────────────────────

def pdf_to_pptx(pdf_path: Path, out_path: Path) -> dict:
    """Convert a PDF to PPTX with real table objects.

    Args:
        pdf_path: Path to the source PDF.
        out_path: Path where the ``.pptx`` file will be saved.

    Returns:
        Summary dict with keys ``tables``, ``text_slides``, ``total_slides``,
        ``size_kb``.
    """
    tables = _extract_tables_from_pdf(pdf_path)
    text_by_page = _extract_text_by_page(pdf_path)

    table_pages = set()
    for t in tables:
        for p in range(t["page"], t["end_page"] + 1):
            table_pages.add(p)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    n_table_slides = 0
    n_text_slides = 0

    for t in tables:
        header, body = normalize_table(t["rows"])
        if not header:
            continue
        pages_str = (
            f"p.{t['page']}-{t['end_page']}"
            if t["page"] != t["end_page"]
            else f"p.{t['page']}"
        )
        slide_title = (
            f"{t['title']}  [{pages_str}]" if t["title"] else f"테이블  [{pages_str}]"
        )
        _add_table_slide(prs, slide_title, header, body)
        n_table_slides += 1

    for page_no in sorted(text_by_page):
        if page_no in table_pages:
            continue
        text = text_by_page[page_no].strip()
        if text:
            _add_text_slide(prs, page_no, text)
            n_text_slides += 1

    prs.save(str(out_path))
    return {
        "tables": n_table_slides,
        "text_slides": n_text_slides,
        "total_slides": n_table_slides + n_text_slides,
        "size_kb": out_path.stat().st_size // 1024,
    }


def main() -> None:
    targets = [
        Path("테스트문서/병합된 테이블 테스트 샘플.pdf"),
        Path("테스트문서/복잡한_테이블_테스트_샘플.pdf"),
    ]

    for pdf_path in targets:
        out_path = pdf_path.with_suffix(".pptx")
        print(f"\n변환 중: {pdf_path.name}")
        info = pdf_to_pptx(pdf_path, out_path)
        print(f"  테이블 슬라이드  : {info['tables']}개")
        print(f"  텍스트 슬라이드  : {info['text_slides']}개")
        print(f"  총 슬라이드      : {info['total_slides']}개")
        print(f"  파일 크기        : {info['size_kb']:,} KB")
        print(f"  저장 경로        : {out_path}")


if __name__ == "__main__":
    sys.stdout.reconfigure(encoding="utf-8")
    main()

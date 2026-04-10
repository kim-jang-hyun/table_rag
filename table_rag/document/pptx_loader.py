"""PPTX document loading and text/table chunking.

Requires: pip install python-pptx
"""

from pathlib import Path
from typing import List

from ..config import DEFAULT_CHUNK_OVERLAP, DEFAULT_CHUNK_SIZE
from ..models import Chunk
from ..table.normalizer import normalize_cell
from ..table.serializer import table_to_text
from .pdf_loader import chunk_text


def load_pptx_chunks(
    pptx_path: Path,
    *,
    extract_table_chunks: bool = True,
    chunk_size: int = DEFAULT_CHUNK_SIZE,
    chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
) -> List[Chunk]:
    """Extract text and table chunks from a PPTX file using python-pptx.

    Args:
        pptx_path:            Path to the ``.pptx`` / ``.ppt`` file.
        extract_table_chunks: When ``True``, table shapes are serialized as
                              separate ``"table"`` chunks.
        chunk_size:           Maximum character length of each text chunk.
        chunk_overlap:        Character overlap between consecutive text chunks.

    Returns:
        List of :class:`~table_rag.models.Chunk` objects.  Each chunk has
        ``chunk_id`` in the form ``p<slide>_t<n>`` (text) or ``p<slide>_tb<n>``
        (table), without a doc-slug prefix.

    Raises:
        RuntimeError: If python-pptx is not installed.
    """
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx is not installed.  Run `pip install python-pptx`."
        ) from exc

    prs = Presentation(str(pptx_path))
    chunks: List[Chunk] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts: List[str] = []

        for shape in slide.shapes:
            if extract_table_chunks and shape.has_table:
                rows: List[List[str]] = [
                    [normalize_cell(cell.text) for cell in row.cells]
                    for row in shape.table.rows
                ]
                table_text = table_to_text(rows, start_page=slide_idx, end_page=slide_idx)
                if table_text:
                    chunks.append(
                        Chunk(
                            chunk_id=f"p{slide_idx:04d}_tb{len(chunks) + 1:03d}",
                            page=slide_idx,
                            text=table_text,
                            source_type="table",
                        )
                    )
            elif shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    slide_texts.append(t)

        full_text = " ".join(" ".join(slide_texts).split())
        if full_text:
            for j, piece in enumerate(
                chunk_text(full_text, chunk_size=chunk_size, chunk_overlap=chunk_overlap),
                start=1,
            ):
                if piece:
                    chunks.append(
                        Chunk(
                            chunk_id=f"p{slide_idx:04d}_t{j:03d}",
                            page=slide_idx,
                            text=piece,
                            source_type="text",
                        )
                    )

    return chunks
